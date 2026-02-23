#!/usr/bin/env python3
"""
Generate PowerPoint Presentation for RWF Founders
Based on PRESENTATION.md content with embedded figures and key metrics.

Usage:
    python generate_presentation.py

Output:
    Presentation_Founders.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Configuration
FIGURES_DIR = "data/results/figures"
OUTPUT_FILE = "Presentation_Founders.pptx"

# Color scheme (professional blue/gray)
PRIMARY_COLOR = RGBColor(0x1a, 0x56, 0x8c)  # Deep blue
SECONDARY_COLOR = RGBColor(0x4a, 0x90, 0xd9)  # Light blue
ACCENT_COLOR = RGBColor(0x2e, 0x7d, 0x32)  # Green for positive
TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)  # Dark gray
LIGHT_GRAY = RGBColor(0x95, 0x95, 0x95)


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, title):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2), Inches(10), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, nudge=None, image_path=None):
    """Add a content slide with title, bullets, optional nudge, and optional image."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # Calculate content area based on image presence
    if image_path and os.path.exists(image_path):
        content_width = 4.5
        content_left = 0.5
        # Add image on right side
        slide.shapes.add_picture(image_path, Inches(5.2), Inches(1.3), width=Inches(4.5))
    else:
        content_width = 9
        content_left = 0.5

    # Bullets
    bullet_top = 1.3
    bullet_box = slide.shapes.add_textbox(Inches(content_left), Inches(bullet_top),
                                          Inches(content_width), Inches(4))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(12)

    # Nudge (friendly explanation) at bottom
    if nudge:
        nudge_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.6))
        tf = nudge_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"💡 {nudge}"
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = SECONDARY_COLOR

    return slide


def add_table_slide(prs, title, headers, rows, nudge=None):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header

    table_width = min(9, num_cols * 2.2)
    col_width = table_width / num_cols

    table = slide.shapes.add_table(num_rows, num_cols,
                                   Inches(0.5), Inches(1.3),
                                   Inches(table_width), Inches(0.5 * num_rows)).table

    # Set column widths
    for i in range(num_cols):
        table.columns[i].width = Inches(col_width)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_COLOR
            p.alignment = PP_ALIGN.CENTER
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xf5, 0xf5, 0xf5)

    # Nudge
    if nudge:
        nudge_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.6))
        tf = nudge_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"💡 {nudge}"
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = SECONDARY_COLOR

    return slide


def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets, nudge=None):
    """Add a slide with two columns."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4.3), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR

    # Left bullets
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.3), Inches(3))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR

    # Right bullets
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(3))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)

    # Nudge
    if nudge:
        nudge_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.6))
        tf = nudge_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"💡 {nudge}"
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = SECONDARY_COLOR

    return slide


def add_image_slide(prs, title, image_path, caption=None, nudge=None):
    """Add a slide with a large centered image."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # Image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.2), width=Inches(8))
    else:
        # Placeholder text if image not found
        placeholder = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        tf = placeholder.text_frame
        p = tf.paragraphs[0]
        p.text = f"[Image not found: {image_path}]"
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.4))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Nudge
    if nudge:
        nudge_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.1), Inches(9), Inches(0.5))
        tf = nudge_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"💡 {nudge}"
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = SECONDARY_COLOR

    return slide


def add_dual_image_slide(prs, title, left_image, right_image, left_caption=None, right_caption=None, nudge=None):
    """Add a slide with two images side by side."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # Left image
    if os.path.exists(left_image):
        slide.shapes.add_picture(left_image, Inches(0.3), Inches(1.2), width=Inches(4.6))

    # Right image
    if os.path.exists(right_image):
        slide.shapes.add_picture(right_image, Inches(5.1), Inches(1.2), width=Inches(4.6))

    # Captions
    if left_caption:
        cap_box = slide.shapes.add_textbox(Inches(0.3), Inches(4.5), Inches(4.6), Inches(0.4))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_caption
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

    if right_caption:
        cap_box = slide.shapes.add_textbox(Inches(5.1), Inches(4.5), Inches(4.6), Inches(0.4))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_caption
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Nudge
    if nudge:
        nudge_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.6))
        tf = nudge_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"💡 {nudge}"
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = SECONDARY_COLOR

    return slide


def main():
    """Generate the complete presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

    # ===========================================
    # SLIDE 1: Title
    # ===========================================
    add_title_slide(
        prs,
        "RWF Economic Impact Model",
        "Lifetime Economic Benefits Analysis  •  Updated February 2026"
    )

    # ===========================================
    # SLIDE 2: Executive Summary
    # ===========================================
    add_table_slide(
        prs,
        "Executive Summary: What We Found",
        ["Metric", "RTE", "Apprenticeship"],
        [
            ["Average LNPV", "Rs 14.0 Lakhs", "Rs 34.4 Lakhs"],
            ["Program Cost", "Rs 1.04 Lakhs", "Rs 1.58 Lakhs"],
            ["Average BCR", "13.5 : 1", "21.7 : 1"],
            ["BCR Range", "5.0 - 27.6 : 1", "11.8 - 33.1 : 1"],
            ["All 32 scenarios positive?", "Yes ✓", "Yes ✓"],
        ],
        nudge="BCR = Benefit-Cost Ratio. A BCR of 13.5:1 means every Rs 1 invested generates Rs 13.5 in lifetime benefits."
    )

    # ===========================================
    # SLIDE 3: The Investment Multiplier (Three Legs)
    # ===========================================
    add_table_slide(
        prs,
        "The Investment Multiplier: Three Legs of Value",
        ["Stage", "RTE", "Apprenticeship"],
        [
            ["Leg 1: RWF Direct Spend", "Rs 4,000", "Rs 6,000"],
            ["Leg 2: Unlocked Funds", "Rs 1,04,000 (govt)", "Rs 1,58,460 (govt+private)"],
            ["Leg 3: Lifetime Value", "~Rs 14 Lakhs", "~Rs 34.4 Lakhs"],
            ["Ratio", "1 : 26 : 350", "1 : 26 : 567"],
        ],
        nudge="Every Rs 1 RWF spends unlocks Rs 26 in government/private funds, which generates Rs 350-567 in lifetime value."
    )

    # ===========================================
    # SLIDE 4: India's Labor Market Challenge
    # ===========================================
    add_two_column_slide(
        prs,
        "The Challenge: India's Bifurcated Labor Market",
        "Formal Sector (~10%)",
        [
            "Entry wage: Rs 32,800/month",
            "Annual growth: +1.5%/year",
            "Contracts, PF, ESI, job security",
            "By Year 40: 4× informal wages"
        ],
        "Informal Sector (~90%)",
        [
            "Entry wage: Rs 13,425/month",
            "Annual growth: -0.2%/year",
            "No contracts, no benefits",
            "Wage stagnation over career"
        ],
        nudge="This is why P(Formal) is our #1 NPV driver — shifting someone to formal employment has larger lifetime effects than any initial wage premium."
    )

    # ===========================================
    # SLIDE 5: How RTE Creates Value
    # ===========================================
    add_content_slide(
        prs,
        "How RTE Creates Value",
        [
            "Private School → +0.137 SD test score gains",
            "Educational credentials → signaling to employers",
            "Formal sector entry: 30% vs 9.1% baseline (3.3× improvement)",
            "Mincer wage premium: 5.8% per year of quality education",
            "",
            "Key Finding: 79% of value comes from formal sector access,",
            "only 21% from improved learning outcomes"
        ],
        nudge="The primary value of RTE is as a pathway to formal employment, not just better learning. Career guidance can maximize this effect.",
        image_path=os.path.join(FIGURES_DIR, "decomposition_stacked_bar.png")
    )

    # ===========================================
    # SLIDE 6: RTE Results
    # ===========================================
    add_table_slide(
        prs,
        "RTE Results by Scenario",
        ["Scenario", "LNPV", "BCR"],
        [
            ["South Male Urban (highest)", "Rs 28.7L", "27.6 : 1"],
            ["West Male Urban", "Rs 22.2L", "21.4 : 1"],
            ["North Female Urban", "Rs 11.4L", "11.0 : 1"],
            ["East Female Rural (lowest)", "Rs 5.2L", "5.0 : 1"],
            ["Average (all 16)", "Rs 14.0L", "13.5 : 1"],
        ],
        nudge="Even in the most pessimistic scenario (East Female Rural), RTE delivers a 5:1 return — well above the 3:1 threshold for development programs."
    )

    # ===========================================
    # SLIDE 7: How Apprenticeship Creates Value
    # ===========================================
    add_content_slide(
        prs,
        "How Apprenticeship Creates Value",
        [
            "Program → NAC Certification (85% completion)",
            "Employer Absorption: 68% formal placement (vs 9% baseline)",
            "Initial premium: Rs 78,000/year above counterfactual",
            "Premium decays with 12-year half-life",
            "",
            "Key Finding: The 68% placement rate creates a 59 p.p.",
            "advantage that dominates the calculation, even with decay"
        ],
        nudge="Skill premiums fade over time, but formal sector entry is permanent. That's why Apprenticeship NPV remains high despite decay.",
        image_path=os.path.join(FIGURES_DIR, "validation_decay_trajectory.png")
    )

    # ===========================================
    # SLIDE 8: Apprenticeship Results
    # ===========================================
    add_table_slide(
        prs,
        "Apprenticeship Results by Scenario",
        ["Scenario", "LNPV", "BCR"],
        [
            ["South Male Urban (highest)", "Rs 52.3L", "33.1 : 1"],
            ["West Male Urban", "Rs 50.4L", "31.9 : 1"],
            ["North Female Urban", "Rs 36.0L", "22.8 : 1"],
            ["East Female Rural (lowest)", "Rs 18.6L", "11.8 : 1"],
            ["Average (all 16)", "Rs 34.4L", "21.7 : 1"],
        ],
        nudge="Apprenticeship assumptions are conservative: premium is 33% of theoretical value, half-life is below international benchmarks (15-20 years)."
    )

    # ===========================================
    # SLIDE 9: Regional Analysis
    # ===========================================
    add_image_slide(
        prs,
        "Regional Analysis: Where Returns Are Highest",
        os.path.join(FIGURES_DIR, "boxplot_regional.png"),
        caption="South/West regions show 20-50% higher returns than North/East",
        nudge="Geographic targeting can significantly improve cost-effectiveness. Urban beneficiaries show 30-50% higher LNPV than rural."
    )

    # ===========================================
    # SLIDE 10: Sensitivity Analysis
    # ===========================================
    add_dual_image_slide(
        prs,
        "Sensitivity Analysis: What Parameters Matter Most?",
        os.path.join(FIGURES_DIR, "tornado_rte.png"),
        os.path.join(FIGURES_DIR, "tornado_apprenticeship.png"),
        left_caption="RTE: P(Formal) is the top driver",
        right_caption="Apprenticeship: Placement rate dominates",
        nudge="Tornado diagrams show how LNPV changes when each parameter varies across its uncertainty range. Wider bars = more influential parameters."
    )

    # ===========================================
    # SLIDE 11: Monte Carlo Uncertainty
    # ===========================================
    add_dual_image_slide(
        prs,
        "Uncertainty Quantification: 1,000 Monte Carlo Simulations",
        os.path.join(FIGURES_DIR, "histogram_monte_carlo_rte.png"),
        os.path.join(FIGURES_DIR, "histogram_monte_carlo_apprenticeship.png"),
        left_caption="RTE: Median Rs 13.6L, 90% CI Rs 5.2-27.5L",
        right_caption="Apprenticeship: Median Rs 40.5L, 90% CI Rs 22-62L",
        nudge="Monte Carlo draws random parameter values 1,000 times to show the distribution of possible outcomes. P(positive) = 100% for both interventions."
    )

    # ===========================================
    # SLIDE 12: Cost-Benefit Positioning
    # ===========================================
    add_image_slide(
        prs,
        "Cost-Benefit Positioning: Massive Safety Margin",
        os.path.join(FIGURES_DIR, "breakeven_bar_chart.png"),
        caption="Both interventions operate far below break-even thresholds",
        nudge="Actual costs are 5× below RTE's break-even and 8× below Apprenticeship's break-even at BCR=3:1. Strong investment case."
    )

    # ===========================================
    # SLIDE 13: DUAL BCR ANALYSIS (Added Feb 2026)
    # ===========================================
    add_table_slide(
        prs,
        "Dual BCR Analysis: Two Perspectives on ROI",
        ["Metric", "RTE", "Apprenticeship"],
        [
            ["Average LNPV", "Rs 14.9 Lakhs", "Rs 36.4 Lakhs"],
            ["Full BCR (Total Investment)", "14.3:1 (5.4-29.0)", "23.0:1 (12.5-35.0)"],
            ["RWF-only BCR (Direct Spend)", "372:1 (141-755)", "606:1 (329-924)"],
            ["RWF Direct Cost", "Rs 4,000", "Rs 6,000"],
            ["Total Investment", "Rs 1.04 Lakhs", "Rs 1.58 Lakhs"],
            ["Unlock Multiplier", "26×", "39.6×"],
        ],
        nudge="Full BCR: Total benefits / total costs (RWF + govt + private). RWF-only BCR: Total benefits / RWF direct spend only."
    )

    # ===========================================
    # SLIDE 14: BCR Sensitivity Analysis (Added Feb 2026)
    # ===========================================
    add_table_slide(
        prs,
        "BCR Sensitivity: Discount Rates",
        ["Discount Rate", "RTE LNPV", "RTE BCR (Full)", "App LNPV", "App BCR (Full)"],
        [
            ["3% (low)", "Rs 32.8L", "31.5:1", "Rs 74.9L", "47.2:1"],
            ["5% (central)", "Rs 23.5L", "22.6:1", "Rs 53.4L", "33.7:1"],
            ["8% (high)", "Rs 15.7L", "15.1:1", "Rs 35.1L", "22.2:1"],
        ],
        nudge="Even at 8% discount rate (conservative), BCRs remain highly positive. All scenarios exceed the 3:1 development program threshold."
    )

    # ===========================================
    # SLIDE 15: Gender Funnel Analysis
    # ===========================================
    add_two_column_slide(
        prs,
        "Gender Funnel: Current Data & Key Gaps",
        "Apprenticeship Funnel",
        [
            "Mobilized > Registered > Enrolled > Completed > Certified > Placed > Formal > Retained",
            "P(Formal|Apprentice) = 68% (blended)",
            "Gender breakdown: NOT YET AVAILABLE",
            "Female share est. ~30-40% of enrollees",
            "If female P(Formal) is 55% vs 68%, model overstates female NPV by ~24%",
        ],
        "RTE Funnel",
        [
            "Mobilized > Enrolled > Completed Gr 8 > Completed Gr 12 > Labor Market > Formal",
            "P(Formal|RTE) = 30% (expert assumption)",
            "Gender breakdown: NOT YET AVAILABLE",
            "Female LFPR = 37% vs 78% male (PLFS)",
            "Gender wage gap: 24% formal, 32% informal (urban)",
        ],
        nudge="Immediate action: re-query RWF's 68% placement data by gender. This single data point would improve the model at near-zero cost."
    )

    # ===========================================
    # SLIDE 16: Trade Mix & Half-Life
    # ===========================================
    add_table_slide(
        prs,
        "Trade Mix & Skill Half-Life: NPV by Trade Category",
        ["Category", "Example Trades", "Half-Life", "NPV vs Baseline"],
        [
            ["A: Rapid Obsolescence", "IT/ITES, Digital Marketing", "5-8 yrs", "57-84%"],
            ["B: Moderate Decay", "Auto Mechanic, CNC, COPA", "8-12 yrs", "92-114%"],
            ["C: Durable Skills", "Electrician, Plumber, Welder", "12-18 yrs", "114-145%"],
            ["D: Long-Term Persistent", "Health Asst, Draughtsman", "18-25 yrs", "145-170%"],
        ],
        nudge="National NATS trade mix is dominated by Category C (43% share), suggesting the model baseline (h=10) may be slightly conservative. Weighted avg h ~ 11.4 years."
    )

    # ===========================================
    # SLIDE 17: Trade Mix Scenarios
    # ===========================================
    add_table_slide(
        prs,
        "Trade Mix Scenarios: BCR Impact",
        ["Scenario", "Weighted h", "NPV (approx.)", "Full BCR", "vs Baseline"],
        [
            ["IT-Heavy (pessimistic)", "9.2 yrs", "Rs 7.5L", "4.7:1", "-6%"],
            ["Balanced National", "11.4 yrs", "Rs 8.8L", "5.6:1", "+10%"],
            ["Manufacturing-Heavy", "15.1 yrs", "Rs 10.5L", "6.7:1", "+32%"],
            ["Model baseline", "10 yrs", "Rs 8.0L", "5.1:1", "--"],
        ],
        nudge="Steering youth toward durable trades (electrician, plumber, welder) could increase apprenticeship NPV by ~31%."
    )

    # ===========================================
    # SLIDE 18: Top 5 Evidence Sources
    # ===========================================
    add_table_slide(
        prs,
        "Top 5 Evidence Sources",
        ["#", "Source", "Used For", "Sample", "Confidence"],
        [
            ["1", "PLFS 2023-24", "Baseline wages, employment", "~433K individuals", "HIGH"],
            ["2", "ILO India Employment 2024", "Formal entry rates, youth", "PLFS microdata", "HIGH"],
            ["3", "Chen et al. (2022)", "Mincer returns (5.8%/yr)", "~100K households", "MODERATE"],
            ["4", "Muralidharan (2013)", "RTE test score gain (0.137 SD)", "18,926 obs (RCT)", "MODERATE"],
            ["5", "RWF Data (Nov 2025)", "Placement rate (68%)", "RWF participants", "HIGH"],
        ],
        nudge="Of 12 high-impact parameters, only 1 has HIGH confidence (P_FORMAL_APPRENTICE). The tracer study targets the 3 LOW-confidence parameters."
    )

    # ===========================================
    # SLIDE 19: Tracer Study Plan
    # ===========================================
    add_content_slide(
        prs,
        "Tracer Study: Closing the Evidence Gaps",
        [
            "Sample: 1,200-1,600 (300-400 RTE + 300-400 Apprentice + comparison group)",
            "Cohorts: 2015-2024 (expanded vintage for cross-cohort decay estimation)",
            "",
            "Priority research questions:",
            "  1. P_FORMAL_RTE validation (currently 30%, zero empirical data)",
            "  2. External validation of P_FORMAL_APPRENTICE (currently 68%)",
            "  3. Trade-specific decay estimation (IT vs electrician vs health)",
            "  4. Gender-disaggregated outcomes (placement, wages, retention)",
            "",
            "Timeline: 18 months core  |  Budget: Rs 45-55 lakhs",
            "Expected: reduce parameter uncertainty from +/-50% to +/-15-20%"
        ],
        nudge="One investment (Rs 45-55L) closes 7 of 8 evidence gaps simultaneously. This is the single most cost-effective next step."
    )

    # ===========================================
    # SLIDE 20: Policy Persistence (Illustrative)
    # ===========================================
    add_table_slide(
        prs,
        "Policy Persistence: Multi-Cohort Impact (Illustrative)",
        ["Scenario", "Years Post-Exit", "Attribution", "Multiplier", "Effective BCR (RTE)"],
        [
            ["Conservative", "5 yrs", "20%", "1.86x", "10.2x"],
            ["Moderate", "10 yrs", "40%", "4.40x", "62.9x"],
            ["Optimistic", "15 yrs", "60%", "9.31x", "395.0x"],
        ],
        nudge="ILLUSTRATIVE ONLY: Attribution fractions have no empirical basis. Shows how persistence compounds value — single-cohort BCR is a conservative floor."
    )

    # ===========================================
    # SLIDE 21: Ask & Next Steps
    # ===========================================
    add_content_slide(
        prs,
        "Ask & Next Steps",
        [
            "1. Approve tracer study launch (Rs 45-55L, 18 months)",
            "   -> Validates the 3 LOW-confidence parameters that drive BCR",
            "",
            "2. Request actual cost data (RWF direct spend per beneficiary)",
            "   -> Enables precise RWF-only BCR calculation",
            "",
            "3. Re-query 68% placement rate by gender (1-2 weeks, near-zero cost)",
            "   -> Immediate model improvement",
            "",
            "4. Microdata access: PLFS unit records + RWF beneficiary contact database",
            "   -> Needed for tracer sampling frame and enhanced calibration",
            "",
            "5. Pilot lightweight exit survey for next 2-3 completing cohorts",
            "   -> Generates gender-disaggregated data before tracer results arrive"
        ],
        nudge=None
    )

    # ===========================================
    # SLIDE 22: Decision Framework
    # ===========================================
    add_table_slide(
        prs,
        "Decision Framework: Which Intervention When?",
        ["If your priority is...", "Consider...", "Because..."],
        [
            ["Max per-beneficiary impact", "Apprenticeship", "2.5× higher LNPV"],
            ["Max reach with limited budget", "RTE", "Lower cost, simpler delivery"],
            ["Serve underserved regions", "Targeted Apprenticeship", "Higher marginal returns"],
            ["Long-term systemic change", "RTE", "Generational education shift"],
            ["Quick wins / demonstrable outcomes", "Apprenticeship", "Faster employment results"],
        ],
        nudge=None
    )

    # ===========================================
    # APPENDIX: Key Assumptions & Uncertainty
    # ===========================================
    add_section_slide(prs, "Appendix")

    add_table_slide(
        prs,
        "Key Assumptions & Confidence Levels",
        ["Parameter", "Value", "Range", "Confidence", "Tracer?"],
        [
            ["P(Formal|RTE)", "30%", "20-50%", "LOW", "Yes"],
            ["P(Formal|Apprentice)", "68%", "50-90%", "HIGH", "Yes"],
            ["Decay Half-Life", "10 yrs", "5-50 yrs", "LOW", "Partial"],
            ["RTE Retention", "60%", "50-75%", "LOW", "Yes"],
            ["Completion Rate", "85%", "75-95%", "LOW", "Yes"],
            ["Mincer Return", "5.8%/yr", "5-8%", "MODERATE", "No"],
            ["Discount Rate", "5%", "3-8%", "MODERATE", "No"],
            ["Wage Growth (Formal)", "1.5%/yr", "0.5-2.5%", "MODERATE", "No"],
        ],
        nudge="Of 12 high-impact parameters, only 1 has HIGH confidence. The tracer study (Rs 45-55L) would validate the 4 most uncertain parameters."
    )

    # ===========================================
    # CLOSING SLIDE
    # ===========================================
    add_title_slide(
        prs,
        "Thank You",
        "Both interventions generate positive returns in ALL 32 scenarios tested.\n"
        "The key question is how to optimize delivery for maximum impact.\n\n"
        "RWF Economic Impact Model  •  February 2026"
    )

    # Save presentation
    prs.save(OUTPUT_FILE)
    print(f"✓ Presentation saved to: {OUTPUT_FILE}")
    print(f"  - Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
